"""
make_ppt.py — CeO2 합성 논문 파이프라인 발표 자료 생성
실행: python make_ppt.py
출력: output/ceria_pipeline_presentation.pptx
"""
import os, json, glob
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import font_manager
import warnings
warnings.filterwarnings("ignore")

# ── 경로 설정 ──────────────────────────────────────────────────────────────────
BASE_DIR  = os.path.dirname(os.path.abspath(__file__))
OUT_DIR   = os.path.join(BASE_DIR, "output")
MODEL_DIR = os.path.join(OUT_DIR, "model")
PPT_PATH  = os.path.join(OUT_DIR, "ceria_pipeline_presentation.pptx")
IMG_DIR   = os.path.join(OUT_DIR, "ppt_images")
os.makedirs(IMG_DIR, exist_ok=True)

# 한글 폰트 설정
_avail = {f.name for f in font_manager.fontManager.ttflist}
for _fn in ["Malgun Gothic", "NanumGothic", "AppleGothic", "DejaVu Sans"]:
    if _fn in _avail:
        plt.rcParams["font.family"] = _fn
        break
plt.rcParams["axes.unicode_minus"] = False

# ── 색상 팔레트 ────────────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1A, 0x3A, 0x5C)   # 타이틀 배경
BLUE_MID   = RGBColor(0x1F, 0x6B, 0xA8)   # 서브 타이틀
BLUE_LIGHT = RGBColor(0xD6, 0xE8, 0xF7)   # 배경 박스
ACCENT     = RGBColor(0xE8, 0x6A, 0x1A)   # 강조 (주황)
GREEN      = RGBColor(0x1E, 0x8A, 0x44)   # 긍정 지표
GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)   # 슬라이드 배경
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_TEXT  = RGBColor(0x55, 0x55, 0x55)

# ── 데이터 로드 ────────────────────────────────────────────────────────────────
def load_data():
    d = {}
    # Excel DB
    xlsx = os.path.join(OUT_DIR, "ceria_synthesis_database.xlsx")
    if os.path.exists(xlsx):
        df = pd.read_excel(xlsx, sheet_name=0)
        if "doi" not in df.columns:
            raw = pd.read_excel(xlsx, sheet_name=0, header=None, nrows=15)
            for idx, row in raw.iterrows():
                if any(str(v).strip().lower() == "doi" for v in row):
                    df = pd.read_excel(xlsx, sheet_name=0, header=idx)
                    break
        d["df"] = df
        d["n_papers"] = len(df)
        d["n_ml"] = int(df["completeness_score"].ge(40).sum()) if "completeness_score" in df.columns else 238
        if "synthesis_method" in df.columns:
            d["method_dist"] = df["synthesis_method"].dropna().value_counts().head(8)
        if "year" in df.columns:
            d["year_dist"] = df["year"].dropna().value_counts().sort_index()

    # CSV 샘플
    csv_path = os.path.join(OUT_DIR, "ceria_samples_merged.csv")
    if os.path.exists(csv_path):
        csv = pd.read_csv(csv_path, low_memory=False)
        d["n_samples"] = len(csv)
        if "particle_size_primary_nm" in csv.columns:
            valid = pd.to_numeric(csv["particle_size_primary_nm"], errors="coerce").dropna()
            d["n_size"] = len(valid)
            d["size_coverage"] = len(valid) / len(csv) * 100
            d["size_median"] = float(valid.median())
            d["size_mean"]   = float(valid.mean())

    # text/ 파일 수
    d["n_fulltext"] = len(glob.glob(os.path.join(BASE_DIR, "text", "*.txt")))

    # ML 피처 중요도
    imp_path = os.path.join(MODEL_DIR, "importance_particle_size_primary_nm.png")
    d["importance_img"] = imp_path if os.path.exists(imp_path) else None

    # 능동학습 HistGBM
    al_path = os.path.join(MODEL_DIR, "active_learning_size_histgbm.csv")
    if os.path.exists(al_path):
        d["active_learning"] = pd.read_csv(al_path)

    # DKL-GP 능동학습
    dkl_al = os.path.join(MODEL_DIR, "dkl_active_learning_particle_size_primary_nm.csv")
    if os.path.exists(dkl_al):
        d["dkl_active"] = pd.read_csv(dkl_al)

    # 역설계
    inv_cube = os.path.join(MODEL_DIR, "inverse_design_cube_10nm.csv")
    inv_rod  = os.path.join(MODEL_DIR, "inverse_design_rod_30nm.csv")
    if os.path.exists(inv_cube):
        d["inv_cube"] = pd.read_csv(inv_cube)
    if os.path.exists(inv_rod):
        d["inv_rod"] = pd.read_csv(inv_rod)

    return d


# ── 그래프 생성 헬퍼 ──────────────────────────────────────────────────────────
def make_year_chart(year_dist, path):
    fig, ax = plt.subplots(figsize=(10, 4))
    years = year_dist.index.astype(int)
    counts = year_dist.values
    bars = ax.bar(years, counts, color="#1F6BA8", edgecolor="white", linewidth=0.5)
    ax.set_xlabel("연도", fontsize=12)
    ax.set_ylabel("논문 수", fontsize=12)
    ax.set_title("연도별 CeO₂ 합성 논문 수 (1990–2026)", fontsize=14, fontweight="bold")
    ax.set_xlim(1988, 2027)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(axis="x", labelsize=9)
    # 최근 3년 강조
    for bar, yr in zip(bars, years):
        if yr >= 2024:
            bar.set_color("#E86A1A")
    plt.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)

def make_method_chart(method_dist, path):
    fig, ax = plt.subplots(figsize=(8, 5))
    colors = ["#1F6BA8","#2E86C1","#3498DB","#5DADE2","#85C1E9",
              "#AED6F1","#D6EBF2","#E8F4FD"]
    methods = [m[:20] for m in method_dist.index]
    bars = ax.barh(methods[::-1], method_dist.values[::-1],
                   color=colors[:len(method_dist)], edgecolor="white")
    for bar, val in zip(bars, method_dist.values[::-1]):
        ax.text(bar.get_width() + 5, bar.get_y() + bar.get_height()/2,
                f"{val:,}", va="center", fontsize=10, color="#333333")
    ax.set_xlabel("논문 수", fontsize=12)
    ax.set_title("합성 방법별 분포 (상위 8개)", fontsize=14, fontweight="bold")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.set_xlim(0, method_dist.values.max() * 1.2)
    plt.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)

def make_pipeline_diagram(path):
    fig, ax = plt.subplots(figsize=(12, 3.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 4)
    ax.axis("off")

    stages = [
        ("Stage 0\n논문 수집", 0.6, "#1A3A5C", "0_collect.py\n7,278편"),
        ("Stage 1\n전문 수집", 2.5, "#1F6BA8", "1_download.py\n5,426편 (74.5%)"),
        ("Stage 2\nGPT 추출", 4.4, "#2980B9", "2~5번\n6,236 샘플"),
        ("Stage 3\n후처리", 6.3, "#3498DB", "6~11번\nExcel/JSONL"),
        ("Stage 4\nML 학습", 8.2, "#5DADE2", "12_model.py\nR²=-0.038"),
        ("DKL-GP\n불확실성", 10.1, "#E86A1A", "12c_gpr.py\nlog-R²=0.307"),
    ]

    for label, x, color, sub in stages:
        rect = mpatches.FancyBboxPatch((x-0.75, 0.8), 1.5, 2.0,
            boxstyle="round,pad=0.1", facecolor=color, edgecolor="white", linewidth=2)
        ax.add_patch(rect)
        ax.text(x, 2.15, label, ha="center", va="center",
                fontsize=9, fontweight="bold", color="white",
                multialignment="center")
        ax.text(x, 1.2, sub, ha="center", va="center",
                fontsize=8, color="white", alpha=0.9, multialignment="center")

    # 화살표
    for i in range(len(stages)-1):
        x1 = stages[i][1] + 0.75
        x2 = stages[i+1][1] - 0.75
        ax.annotate("", xy=(x2, 1.8), xytext=(x1, 1.8),
                    arrowprops=dict(arrowstyle="->", color="#555555", lw=2))

    ax.set_title("파이프라인 아키텍처", fontsize=14, fontweight="bold", pad=10)
    plt.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)

def make_ml_results_chart(path):
    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))

    # 왼쪽: R² 비교
    ax = axes[0]
    models = ["HistGBM\n(primary_nm)", "HistGBM\n(TEM)", "HistGBM\nsol-gel", "DKL-GP\n(log scale)"]
    r2_vals = [-0.038, -0.061, 0.111, 0.307]
    colors_r2 = ["#E74C3C" if v < 0 else "#1E8A44" for v in r2_vals]
    bars = ax.bar(models, r2_vals, color=colors_r2, edgecolor="white", linewidth=0.8)
    ax.axhline(0, color="black", linewidth=1, linestyle="--", alpha=0.5)
    for bar, val in zip(bars, r2_vals):
        ypos = val + 0.01 if val >= 0 else val - 0.025
        ax.text(bar.get_x() + bar.get_width()/2, ypos, f"{val:+.3f}",
                ha="center", va="bottom" if val >= 0 else "top", fontsize=10, fontweight="bold")
    ax.set_ylabel("R²", fontsize=12)
    ax.set_title("모델별 R² 성능", fontsize=13, fontweight="bold")
    ax.set_ylim(-0.15, 0.40)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    # 오른쪽: 피처 중요도
    ax2 = axes[1]
    features = ["synthesis_method", "solvent", "ce_precursor",
                "capping_agent", "synth_temp", "mineralizer",
                "synth_time", "anion_type"]
    importances = [22.9, 11.9, 10.6, 6.8, 5.7, 5.2, 4.8, 4.1]
    colors_fi = ["#1F6BA8"] * len(features)
    colors_fi[0] = "#E86A1A"
    y_pos = range(len(features))
    ax2.barh(list(y_pos), importances, color=colors_fi, edgecolor="white")
    ax2.set_yticks(list(y_pos))
    ax2.set_yticklabels(features, fontsize=9)
    ax2.set_xlabel("피처 중요도 (%)", fontsize=11)
    ax2.set_title("HistGBM 피처 중요도 상위 8개", fontsize=13, fontweight="bold")
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_visible(False)
    for i, val in enumerate(importances):
        ax2.text(val + 0.3, i, f"{val:.1f}%", va="center", fontsize=9)
    ax2.set_xlim(0, 28)

    plt.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)

def make_coverage_chart(n_samples, n_size, path):
    fig, ax = plt.subplots(figsize=(7, 4))
    categories = ["전문 보유\n(7,278편 중)", "GPT 추출\n완료", "입자크기\n커버리지\n(TEM+SEM)", "ML 학습\n가능\n(완성도≥40%)"]
    values = [74.5, 59.4, 37.8, 3.3]  # %
    colors = ["#1F6BA8", "#2980B9", "#E86A1A", "#1E8A44"]
    bars = ax.bar(categories, values, color=colors, edgecolor="white", linewidth=0.8, width=0.5)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                f"{val:.1f}%", ha="center", va="bottom", fontsize=11, fontweight="bold")
    ax.set_ylabel("비율 (%)", fontsize=12)
    ax.set_title("데이터 파이프라인 단계별 커버리지", fontsize=14, fontweight="bold")
    ax.set_ylim(0, 90)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.0f}%"))
    plt.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)

def make_size_dist_chart(csv_path, path):
    csv = pd.read_csv(csv_path, low_memory=False)
    col = "particle_size_primary_nm"
    if col not in csv.columns:
        return
    vals = pd.to_numeric(csv[col], errors="coerce").dropna()
    vals = vals[(vals > 0) & (vals < 300)]

    fig, ax = plt.subplots(figsize=(8, 4))
    ax.hist(vals, bins=50, color="#1F6BA8", edgecolor="white", alpha=0.85)
    ax.axvline(vals.median(), color="#E86A1A", linewidth=2, linestyle="--",
               label=f"중앙값 {vals.median():.1f}nm")
    ax.axvline(vals.mean(), color="#1E8A44", linewidth=2, linestyle="-.",
               label=f"평균 {vals.mean():.1f}nm")
    ax.set_xlabel("1차 입자 크기 (nm)", fontsize=12)
    ax.set_ylabel("빈도", fontsize=12)
    ax.set_title(f"1차 입자 크기 분포 (TEM+SEM, n={len(vals):,})", fontsize=14, fontweight="bold")
    ax.legend(fontsize=11)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ── PPT 유틸 ──────────────────────────────────────────────────────────────────
def add_bg(slide, color=GRAY_LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title, subtitle=None):
    # 상단 파란 바
    bar = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(33.87), Cm(3.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE_DARK
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "맑은 고딕"

    tf.margin_left  = Cm(0.8)
    tf.margin_top   = Cm(0.5)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Cm(0.8), Cm(2.6), Cm(30), Cm(0.8))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(11)
        run2.font.color.rgb = GRAY_TEXT
        run2.font.name = "맑은 고딕"

def add_text_box(slide, text, left, top, width, height,
                 font_size=13, bold=False, color=BLACK,
                 bg_color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    if bg_color:
        box.fill.solid(); box.fill.fore_color.rgb = bg_color
        box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return box

def add_kpi_box(slide, label, value, unit, left, top, color=BLUE_MID):
    box = slide.shapes.add_shape(1, Cm(left), Cm(top), Cm(5.5), Cm(3.2))
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    # 값
    tb = slide.shapes.add_textbox(Cm(left+0.2), Cm(top+0.2), Cm(5.1), Cm(1.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = value
    run.font.size = Pt(28); run.font.bold = True
    run.font.color.rgb = WHITE; run.font.name = "맑은 고딕"
    # 단위
    tb2 = slide.shapes.add_textbox(Cm(left+0.2), Cm(top+1.9), Cm(5.1), Cm(0.7))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run(); run2.text = unit
    run2.font.size = Pt(9); run2.font.color.rgb = WHITE; run2.font.name = "맑은 고딕"
    # 라벨
    tb3 = slide.shapes.add_textbox(Cm(left+0.2), Cm(top+2.6), Cm(5.1), Cm(0.6))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run(); run3.text = label
    run3.font.size = Pt(10); run3.font.bold = True
    run3.font.color.rgb = WHITE; run3.font.name = "맑은 고딕"

def add_image(slide, img_path, left, top, width, height=None):
    if not os.path.exists(img_path):
        return
    if height:
        slide.shapes.add_picture(img_path, Cm(left), Cm(top), Cm(width), Cm(height))
    else:
        slide.shapes.add_picture(img_path, Cm(left), Cm(top), Cm(width))

def add_bullet_list(slide, items, left, top, width, height,
                    font_size=12, title=None, title_color=BLUE_MID):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(item, tuple):
            indent, text = item
        else:
            indent, text = 0, item
        p.level = indent
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size - indent)
        run.font.name = "맑은 고딕"
        if indent == 0:
            run.font.bold = True
            run.font.color.rgb = title_color if i == 0 and title else BLACK
        else:
            run.font.color.rgb = GRAY_TEXT


# ── 슬라이드 생성 ─────────────────────────────────────────────────────────────
def make_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 전체 배경
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = BLUE_DARK

    # 타이틀
    tb = slide.shapes.add_textbox(Cm(2), Cm(4), Cm(29), Cm(4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "CeO₂ 나노입자 합성 논문 ML 파이프라인"
    run.font.size = Pt(36); run.font.bold = True
    run.font.color.rgb = WHITE; run.font.name = "맑은 고딕"

    # 서브타이틀
    tb2 = slide.shapes.add_textbox(Cm(2), Cm(8.5), Cm(29), Cm(2))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "7,278편 논문 자동 추출 및 ML 기반 합성 조건 역설계"
    run2.font.size = Pt(18); run2.font.color.rgb = RGBColor(0xAE, 0xD6, 0xF1)
    run2.font.name = "맑은 고딕"

    # 날짜
    tb3 = slide.shapes.add_textbox(Cm(2), Cm(11), Cm(29), Cm(1))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "2026년 6월  |  세리아 합성 AI 연구"
    run3.font.size = Pt(13); run3.font.color.rgb = RGBColor(0x85, 0xC1, 0xE9)
    run3.font.name = "맑은 고딕"

    # 하단 KPI 요약
    n = data.get("n_papers", 7278)
    ft = data.get("n_fulltext", 5426)
    ns = data.get("n_samples", 6236)
    kpis = [
        ("논문", f"{n:,}", "편"),
        ("전문 보유", f"{ft:,}", f"편 ({ft/n*100:.0f}%)"),
        ("추출 샘플", f"{ns:,}", "행"),
        ("ML 모델", "2종", "HistGBM + DKL-GP"),
    ]
    start_x = 2.5
    for label, val, unit in kpis:
        box = slide.shapes.add_shape(1, Cm(start_x), Cm(13.5), Cm(6.5), Cm(2.5))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        box.line.color.rgb = RGBColor(0x5D, 0xAD, 0xE2); box.line.width = Pt(1)

        tb_v = slide.shapes.add_textbox(Cm(start_x+0.2), Cm(13.6), Cm(6.1), Cm(1.2))
        tf_v = tb_v.text_frame
        pv = tf_v.paragraphs[0]; pv.alignment = PP_ALIGN.CENTER
        rv = pv.add_run(); rv.text = val
        rv.font.size = Pt(22); rv.font.bold = True
        rv.font.color.rgb = WHITE; rv.font.name = "맑은 고딕"

        tb_u = slide.shapes.add_textbox(Cm(start_x+0.2), Cm(14.7), Cm(6.1), Cm(0.7))
        tf_u = tb_u.text_frame
        pu = tf_u.paragraphs[0]; pu.alignment = PP_ALIGN.CENTER
        ru = pu.add_run(); ru.text = f"{unit}  |  {label}"
        ru.font.size = Pt(9); ru.font.color.rgb = RGBColor(0xAE, 0xD6, 0xF1)
        ru.font.name = "맑은 고딕"
        start_x += 7.0


def make_overview_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_LIGHT)
    add_title_bar(slide, "프로젝트 개요", "CeO₂ 나노입자 합성 조건 자동 추출 및 ML 역설계")

    items_left = [
        "🎯  연구 목적",
        (1, "CeO₂ 나노입자 합성 논문 7,278편에서"),
        (1, "합성 조건 + 측정 결과 자동 추출"),
        (1, "ML 학습 데이터셋 구축 및 역설계"),
        "",
        "📌  핵심 질문",
        (1, "어떤 합성 조건에서 원하는 입자 크기/형태가 만들어지는가?"),
        (1, "모르는 조합 중 가장 탐색 가치 있는 조건은?"),
        "",
        "🔬  추출 필드 (8개)",
        (1, "합성법, Ce 전구체, 용매, 온도, pH,"),
        (1, "Ce 농도, 광화제 농도, 합성 부피"),
    ]
    add_bullet_list(slide, items_left, 1.0, 3.3, 15.5, 14, font_size=12)

    items_right = [
        "📊  현재 데이터 현황",
        (1, f"논문 DB: 7,278편 (1990~2026)"),
        (1, f"전문 보유: 5,426편 (74.5%)"),
        (1, f"추출 샘플: 6,236행"),
        (1, f"입자크기 커버리지: 37.8% (TEM+SEM)"),
        "",
        "🤖  ML 모델",
        (1, "HistGBM: 29개 피처, GroupKFold(5)"),
        (1, "DKL-GP: 불확실성 정량화 (log-R²=0.307)"),
        (1, "능동학습: 실험 제안 자동화"),
        (1, "역설계: 목표 입자 크기/형태 조건 탐색"),
    ]
    add_bullet_list(slide, items_right, 17.5, 3.3, 15.0, 14, font_size=12)

    # 구분선
    line = slide.shapes.add_shape(1, Cm(16.5), Cm(3.5), Cm(0.08), Cm(13.5))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    line.line.fill.background()


def make_data_collection_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_LIGHT)
    add_title_bar(slide, "데이터 수집 현황", "OpenAlex 다층 쿼리 + 1_download.py 전문 수집")

    # KPI 박스 4개
    n = data.get("n_papers", 7278)
    ft = data.get("n_fulltext", 5426)
    kpis = [
        ("총 논문", "7,278", "편  (1990~2026)", BLUE_DARK),
        ("전문 보유", "5,426", f"편  ({ft/n*100:.1f}%)", BLUE_MID),
        ("PDF 파일", "4,161", "개  (pdf/ 폴더)", RGBColor(0x1A, 0x7A, 0x4A)),
        ("GPT 추출", "4,327", "편  (추출 완료)", ACCENT),
    ]
    x_start = 1.2
    for label, val, unit, color in kpis:
        add_kpi_box(slide, label, val, unit, x_start, 3.5, color)
        x_start += 7.5

    # 연도별 차트
    if "year_dist" in data:
        year_img = os.path.join(IMG_DIR, "year_dist.png")
        make_year_chart(data["year_dist"], year_img)
        add_image(slide, year_img, 1.0, 7.2, 20, 9.5)

    # 수집 전략
    add_text_box(slide, "📥  수집 전략 (1_download.py)", 22.5, 7.2, 10.5, 1.0,
                 font_size=11, bold=True, color=BLUE_MID)
    strategies = [
        "1. 기존 PDF → 텍스트 추출",
        "2. Unpaywall OA URL",
        "3. PMC (NCBI 무료 API)  → 1,278편",
        "4. Semantic Scholar OA PDF",
        "5. Sci-Hub (--scihub)  → 1,752편",
        "",
        "✅ 합계: +1,242편 (이번 배치)",
        "✅ 누적: 4,173 → 5,426편",
    ]
    y = 8.4
    for s in strategies:
        color = GREEN if s.startswith("✅") else BLACK
        add_text_box(slide, s, 22.5, y, 10.5, 0.65, font_size=10, color=color)
        y += 0.72


def make_extraction_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_LIGHT)
    add_title_bar(slide, "데이터 추출 파이프라인", "GPT-4o-mini 기반 합성 조건 자동 추출")

    # 커버리지 차트
    cov_img = os.path.join(IMG_DIR, "coverage.png")
    n_s = data.get("n_samples", 6236)
    n_sz = data.get("n_size", 2352)
    make_coverage_chart(n_s, n_sz, cov_img)
    add_image(slide, cov_img, 1.0, 3.5, 18, 9)

    # 추출 필드 정리
    add_text_box(slide, "📋  추출 필드 (8개)", 20.0, 3.5, 12.5, 0.9,
                 font_size=12, bold=True, color=BLUE_MID)
    fields = [
        ("synthesis_method",       "합성 방법"),
        ("ce_precursor",           "Ce 전구체"),
        ("solvent",                "용매"),
        ("synthesis_temperature_c","합성 온도 (°C)"),
        ("ph_synthesis",           "pH"),
        ("ce_concentration_M",     "Ce 농도 (M)"),
        ("mineralizer_conc_M",     "광화제 농도 (M)"),
        ("synthesis_volume_mL",    "합성 부피 (mL)  ★신규"),
    ]
    y = 4.6
    for field, desc in fields:
        color = ACCENT if "신규" in desc else GRAY_TEXT
        add_text_box(slide, f"• {field}", 20.0, y, 7.5, 0.6, font_size=9, color=BLACK)
        add_text_box(slide, desc, 27.0, y, 5.5, 0.6, font_size=9, color=color)
        y += 0.72

    # 입자크기 분포 차트
    csv_p = os.path.join(OUT_DIR, "ceria_samples_merged.csv")
    if os.path.exists(csv_p):
        size_img = os.path.join(IMG_DIR, "size_dist.png")
        make_size_dist_chart(csv_p, size_img)
        add_image(slide, size_img, 1.0, 12.8, 18, 5.5)

    add_text_box(slide, f"입자크기 커버리지: {data.get('size_coverage', 37.8):.1f}%  "
                        f"(n={data.get('n_size', 2352):,})  |  "
                        f"중앙값 {data.get('size_median', 11):.1f}nm  "
                        f"평균 {data.get('size_mean', 18):.1f}nm",
                 1.0, 13.0, 18, 0.7, font_size=9, color=GRAY_TEXT)


def make_method_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_LIGHT)
    add_title_bar(slide, "합성 방법 분포", "7,278편 논문 기준 합성 방법별 통계")

    if "method_dist" in data:
        meth_img = os.path.join(IMG_DIR, "method_dist.png")
        make_method_chart(data["method_dist"], meth_img)
        add_image(slide, meth_img, 1.0, 3.5, 21, 12)

    # 핵심 인사이트
    insights = [
        "💡  핵심 인사이트",
        (1, "Hydrothermal이 가장 많은 논문 보유 → 데이터 풍부"),
        (1, "Sol-gel: per-method 모델 R²=+0.111 (유일하게 양수)"),
        (1, "Combustion: 빠른 합성, 입자 크기 분산 큼"),
        (1, "Solvothermal: 다양한 형태 제어 가능"),
        "",
        "📈  합성법별 ML 성능 (R², nm 기준)",
        (1, "Sol-gel:        +0.111  ✅"),
        (1, "Hydrothermal:  -0.058"),
        (1, "Precipitation:  -0.066"),
        (1, "전체 통합:       -0.038"),
    ]
    add_bullet_list(slide, insights, 23.0, 3.5, 10.0, 14, font_size=11)


def make_ml_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_LIGHT)
    add_title_bar(slide, "ML 모델 결과", "HistGBM 29피처 + DKL-GP 불확실성 정량화")

    ml_img = os.path.join(IMG_DIR, "ml_results.png")
    make_ml_results_chart(ml_img)
    add_image(slide, ml_img, 1.0, 3.5, 30, 9.5)

    # 모델 상세
    details = [
        "🔧  모델 구성",
        (1, "피처 29개: 18 수치형 + 11 범주형 (TargetEncoder)"),
        (1, "검증: GroupKFold(n=5) by DOI — 논문 단위 분리"),
        (1, "분위수 회귀(Q10/Q90): 예측 구간 추정"),
        "",
        "📊  평가 지표 (nm 기준)",
        (1, "primary_nm:  R²=-0.038,  MAE=32.68nm"),
        (1, "TEM 단독:    R²=-0.061,  MAE=30.04nm"),
        (1, "DKL-GP:      log-R²=+0.307  ← 로그 스케일"),
        "",
        "⚠️  R² 음수 원인 분석",
        (1, "같은 합성법 내 입자크기 분산이 매우 큼"),
        (1, "Ce 농도/부피 등 미추출 변수가 결정적 영향"),
        (1, "→ synthesis_volume_mL 추출 중 (개선 예정)"),
    ]
    add_bullet_list(slide, details, 1.0, 13.5, 32, 5.5, font_size=10)


def make_active_learning_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_LIGHT)
    add_title_bar(slide, "능동학습 — 다음 실험 제안", "불확실성 최대 조건 탐색 (HistGBM Q10/Q90 + DKL-GP σ)")

    # 설명
    add_text_box(slide,
        "불확실성이 가장 높은 합성 조건 = 실험하면 ML 모델 개선 효과가 최대인 조건",
        1.0, 3.3, 32, 0.8, font_size=12, color=GRAY_TEXT)

    # HistGBM 능동학습 테이블
    add_text_box(slide, "📊 HistGBM Q90-Q10 불확실성 최대 조건 (상위 5개)",
                 1.0, 4.2, 20, 0.7, font_size=11, bold=True, color=BLUE_MID)

    if "active_learning" in data:
        al = data["active_learning"].head(5)
        headers = ["합성법", "불확실성 구간 (nm)", "Q10 예측", "Q90 예측", "온도(°C)"]
        col_w   = [7.0, 4.5, 3.0, 3.0, 2.5]
        x_pos   = [1.0, 8.0, 12.5, 15.5, 18.5]
        y = 5.1
        # 헤더
        for h, xp, w in zip(headers, x_pos, col_w):
            add_text_box(slide, h, xp, y, w, 0.65, font_size=9,
                         bold=True, color=WHITE,
                         bg_color=BLUE_MID, align=PP_ALIGN.CENTER)
        y += 0.7
        for _, row in al.iterrows():
            cols_data = [
                str(row.get("synthesis_method",""))[:25],
                f"{row.get('uncertainty_interval_nm', 0):.1f} nm",
                f"{row.get('predicted_q10_nm', 0):.1f} nm",
                f"{row.get('predicted_q90_nm', 0):.1f} nm",
                f"{row.get('synthesis_temperature_c', 0):.0f}°C",
            ]
            bg = GRAY_LIGHT if int(y*10) % 2 == 0 else WHITE
            for val, xp, w in zip(cols_data, x_pos, col_w):
                add_text_box(slide, val, xp, y, w, 0.65, font_size=8,
                             color=BLACK, align=PP_ALIGN.CENTER)
            y += 0.65

    # DKL-GP 능동학습
    add_text_box(slide, "🤖 DKL-GP σ 불확실성 최대 조건 (상위 5개)",
                 1.0, 10.5, 20, 0.7, font_size=11, bold=True, color=ACCENT)

    if "dkl_active" in data:
        dkl = data["dkl_active"].head(5)
        headers2 = ["합성법", "σ (log)", "예측 크기(nm)", "σ 구간(nm)", "온도(°C)"]
        x_pos2   = [1.0, 8.0, 12.0, 16.0, 20.0]
        col_w2   = [7.0, 4.0, 4.0, 4.0, 2.5]
        y = 11.3
        for h, xp, w in zip(headers2, x_pos2, col_w2):
            add_text_box(slide, h, xp, y, w, 0.65, font_size=9,
                         bold=True, color=WHITE, bg_color=ACCENT, align=PP_ALIGN.CENTER)
        y += 0.7
        for _, row in dkl.iterrows():
            unc_int = row.get("uncertainty_interval_nm", 0)
            cols2 = [
                str(row.get("synthesis_method",""))[:25],
                f"{row.get('uncertainty_sigma', 0):.3f}",
                f"{row.get('predicted_mean_nm', 0):.1f} nm",
                f"±{unc_int/2:.1f} nm",
                f"{row.get('synthesis_temperature_c', 0):.0f}°C",
            ]
            for val, xp, w in zip(cols2, x_pos2, col_w2):
                add_text_box(slide, val, xp, y, w, 0.65, font_size=8,
                             color=BLACK, align=PP_ALIGN.CENTER)
            y += 0.65

    # 오른쪽 설명
    add_text_box(slide, "🎯  능동학습 의의",
                 22.5, 4.2, 10.5, 0.7, font_size=11, bold=True, color=BLUE_MID)
    explain = [
        "• 7,278편 × 조건 조합에서",
        "  불확실성 최대 조건 자동 선별",
        "",
        "• 다음 실험 설계 시",
        "  모델 개선 효과 극대화 가능",
        "",
        "• DKL-GP: 베이즈 최적화 기반",
        "  σ = 예측 신뢰도의 역수",
        "",
        "• HistGBM: Q90-Q10 구간 폭",
        "  = 예측 불확실성 proxy",
    ]
    y = 5.1
    for e in explain:
        add_text_box(slide, e, 22.5, y, 10.5, 0.7, font_size=10, color=GRAY_TEXT)
        y += 0.72


def make_inverse_design_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_LIGHT)
    add_title_bar(slide, "역설계 — 목표 입자 조건 탐색",
                  "합성 조건 공간에서 원하는 크기/형태 조건 자동 탐색")

    add_text_box(slide,
        "목표 형태/크기를 입력하면 → 문헌 분포 기반으로 최적 합성 조건 랭킹 제시",
        1.0, 3.3, 32, 0.7, font_size=11, color=GRAY_TEXT)

    # Cube 10nm 결과
    add_text_box(slide, "🎯  목표: Cube 형태 + 10nm",
                 1.0, 4.2, 15, 0.7, font_size=12, bold=True, color=BLUE_MID)
    if "inv_cube" in data:
        cube = data["inv_cube"].head(5)
        y = 5.0
        for i, (_, row) in enumerate(cube.iterrows()):
            rank_color = ACCENT if i == 0 else BLUE_MID
            add_text_box(slide, f"#{i+1}", 1.0, y, 1.0, 0.7,
                         font_size=11, bold=True, color=rank_color)
            method = str(row.get("synthesis_method",""))[:30]
            size   = row.get("pred_size_nm", 0)
            score  = row.get("score", 0)
            add_text_box(slide, f"{method}  →  {size:.1f}nm  (score={score:.4f})",
                         2.0, y, 14.0, 0.7, font_size=10, color=BLACK)
            y += 0.8

    # Rod 30nm 결과
    add_text_box(slide, "🎯  목표: Rod 형태 + 30nm",
                 17.5, 4.2, 15.5, 0.7, font_size=12, bold=True, color=ACCENT)
    if "inv_rod" in data:
        rod = data["inv_rod"].head(5)
        y = 5.0
        for i, (_, row) in enumerate(rod.iterrows()):
            rank_color = ACCENT if i == 0 else BLUE_MID
            add_text_box(slide, f"#{i+1}", 17.5, y, 1.0, 0.7,
                         font_size=11, bold=True, color=rank_color)
            method = str(row.get("synthesis_method",""))[:30]
            size   = row.get("pred_size_nm", 0)
            score  = row.get("score", 0)
            add_text_box(slide, f"{method}  →  {size:.1f}nm  (score={score:.4f})",
                         18.5, y, 14.0, 0.7, font_size=10, color=BLACK)
            y += 0.8

    # 역설계 원리 설명
    box = slide.shapes.add_shape(1, Cm(1.0), Cm(10.5), Cm(31.5), Cm(7.5))
    box.fill.solid(); box.fill.fore_color.rgb = BLUE_LIGHT
    box.line.color.rgb = RGBColor(0xAE, 0xD6, 0xF1); box.line.width = Pt(1)

    add_text_box(slide, "⚙️  역설계 알고리즘",
                 1.5, 10.8, 15, 0.7, font_size=11, bold=True, color=BLUE_DARK)
    algo_steps = [
        "① 목표 크기(nm) + 형태(cube/rod 등) 설정",
        "② 10,000개 무작위 합성 조건 샘플링",
        "③ HistGBM 형태 예측 확률 × 크기 근접도 → 점수 계산",
        "④ 상위 10개 조건 반환 (문헌 분포 기반, 현실적 조건)",
    ]
    y = 11.7
    for s in algo_steps:
        add_text_box(slide, s, 1.5, y, 31, 0.7, font_size=10, color=BLACK)
        y += 0.8

    add_text_box(slide, "📌  활용 방법: 연구자가 원하는 형태/크기 입력 → "
                        "제안된 조건으로 실험 설계 → 불확실성 높은 조건 우선 탐색",
                 1.5, 14.8, 31, 1.0, font_size=10, color=GRAY_TEXT)


def make_future_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, GRAY_LIGHT)
    add_title_bar(slide, "현황 및 향후 계획", "진행 중인 작업 및 개선 로드맵")

    # 현재 진행 중
    add_text_box(slide, "🔄  현재 진행 중 (2026-06-11)", 1.0, 3.5, 31, 0.8,
                 font_size=13, bold=True, color=BLUE_MID)
    current = [
        ("2_extract.py 실행 중", "1,088편 GPT 추출  (~$0.87,  약 6~8시간 소요)"),
        ("완료 후 자동 실행",     "main.py --from 2  →  Stage 3 후처리 + Stage 4 ML 재학습"),
    ]
    y = 4.5
    for title, desc in current:
        box = slide.shapes.add_shape(1, Cm(1.0), Cm(y), Cm(31.5), Cm(1.3))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
        box.line.color.rgb = BLUE_MID; box.line.width = Pt(1)
        add_text_box(slide, f"⏳  {title}", 1.3, y+0.1, 10, 0.65, font_size=10, bold=True, color=BLUE_DARK)
        add_text_box(slide, desc, 11.5, y+0.1, 21, 0.65, font_size=10, color=BLACK)
        y += 1.5

    # 단기 계획
    add_text_box(slide, "📋  단기 계획 (1~2주)", 1.0, 8.0, 31, 0.8,
                 font_size=13, bold=True, color=ACCENT)
    short_term = [
        ("synthesis_volume_mL 재추출",   "4_extract_targeted.py --reset  →  Ce 농도(mol/L) 파생 가능"),
        ("ML 피처 확장",                  "ce_concentration_M 커버리지 개선 → R² 개선 기대"),
        ("DKL-GP 재학습",                 "신규 데이터 반영, 능동학습 추천 갱신"),
    ]
    y = 9.0
    for title, desc in short_term:
        add_text_box(slide, f"• {title}", 1.3, y, 12, 0.7, font_size=10, bold=True, color=BLACK)
        add_text_box(slide, desc, 13.5, y, 19.5, 0.7, font_size=10, color=GRAY_TEXT)
        y += 0.9

    # 장기 계획
    add_text_box(slide, "🚀  장기 계획", 1.0, 12.5, 31, 0.8,
                 font_size=13, bold=True, color=GREEN)
    long_term = [
        "실험 검증 → 능동학습 예측 조건 실제 합성 후 모델 피드백 루프 구성",
        "주간 자동화 (run_weekly.py) → Task Scheduler 등록, 신규 논문 자동 수집/학습",
        "CMP 특화 분석 → 입자 크기/형태 + CMP 성능 상관관계 규명",
        "도핑 최적화 → dopant 종류·농도 vs 입자 특성 전용 모델",
    ]
    y = 13.5
    for item in long_term:
        add_text_box(slide, f"• {item}", 1.3, y, 31, 0.75, font_size=10, color=BLACK)
        y += 0.85


def make_summary_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = BLUE_DARK

    tb = slide.shapes.add_textbox(Cm(2), Cm(2.5), Cm(29), Cm(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = "성과 요약"
    run.font.size = Pt(30); run.font.bold = True
    run.font.color.rgb = WHITE; run.font.name = "맑은 고딕"

    achievements = [
        ("7,278편", "논문 DB 구축\n(1990~2026)"),
        ("74.5%",   "전문 확보율\n(5,426편)"),
        ("6,236",   "추출 샘플\n(행)"),
        ("0.307",   "DKL-GP\nlog-R²"),
        ("+0.111",  "Sol-gel\nnm-R²"),
    ]
    x = 1.5
    for val, label in achievements:
        box = slide.shapes.add_shape(1, Cm(x), Cm(5.5), Cm(5.6), Cm(4.5))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        box.line.color.rgb = RGBColor(0x5D, 0xAD, 0xE2); box.line.width = Pt(1.5)

        tb_v = slide.shapes.add_textbox(Cm(x+0.2), Cm(5.8), Cm(5.2), Cm(1.8))
        tf_v = tb_v.text_frame
        pv = tf_v.paragraphs[0]; pv.alignment = PP_ALIGN.CENTER
        rv = pv.add_run(); rv.text = val
        rv.font.size = Pt(30); rv.font.bold = True
        rv.font.color.rgb = RGBColor(0xF0, 0xB0, 0x27); rv.font.name = "맑은 고딕"

        tb_l = slide.shapes.add_textbox(Cm(x+0.2), Cm(7.5), Cm(5.2), Cm(1.5))
        tf_l = tb_l.text_frame
        pl = tf_l.paragraphs[0]; pl.alignment = PP_ALIGN.CENTER
        rl = pl.add_run(); rl.text = label
        rl.font.size = Pt(11); rl.font.color.rgb = RGBColor(0xAE, 0xD6, 0xF1)
        rl.font.name = "맑은 고딕"
        x += 6.2

    # 핵심 메시지
    tb_msg = slide.shapes.add_textbox(Cm(2), Cm(11.5), Cm(29), Cm(2.5))
    tf_msg = tb_msg.text_frame; tf_msg.word_wrap = True
    p_msg = tf_msg.paragraphs[0]; p_msg.alignment = PP_ALIGN.CENTER
    r_msg = p_msg.add_run()
    r_msg.text = ("GPT 자동 추출 + ML 역설계로 7,278편 논문에서 최적 CeO₂ 합성 조건을 탐색하는\n"
                  "데이터 기반 소재 설계 파이프라인 구축 완료")
    r_msg.font.size = Pt(14); r_msg.font.color.rgb = RGBColor(0xAE, 0xD6, 0xF1)
    r_msg.font.name = "맑은 고딕"

    tb_q = slide.shapes.add_textbox(Cm(2), Cm(15), Cm(29), Cm(1.5))
    tf_q = tb_q.text_frame
    pq = tf_q.paragraphs[0]; pq.alignment = PP_ALIGN.CENTER
    rq = pq.add_run(); rq.text = "감사합니다"
    rq.font.size = Pt(22); rq.font.bold = True
    rq.font.color.rgb = WHITE; rq.font.name = "맑은 고딕"


# ── 메인 ──────────────────────────────────────────────────────────────────────
def main():
    print("데이터 로드 중...")
    data = load_data()
    print(f"  논문: {data.get('n_papers', 0):,}편  |  전문: {data.get('n_fulltext', 0):,}편  "
          f"|  샘플: {data.get('n_samples', 0):,}행")

    print("그래프 생성 중...")
    if "year_dist" in data:
        make_year_chart(data["year_dist"], os.path.join(IMG_DIR, "year_dist.png"))
    if "method_dist" in data:
        make_method_chart(data["method_dist"], os.path.join(IMG_DIR, "method_dist.png"))
    make_pipeline_diagram(os.path.join(IMG_DIR, "pipeline.png"))
    make_ml_results_chart(os.path.join(IMG_DIR, "ml_results.png"))
    make_coverage_chart(
        data.get("n_samples", 6236),
        data.get("n_size", 2352),
        os.path.join(IMG_DIR, "coverage.png")
    )
    csv_p = os.path.join(OUT_DIR, "ceria_samples_merged.csv")
    if os.path.exists(csv_p):
        make_size_dist_chart(csv_p, os.path.join(IMG_DIR, "size_dist.png"))

    print("PPT 생성 중...")
    prs = Presentation()
    prs.slide_width  = Cm(33.87)   # 16:9 와이드
    prs.slide_height = Cm(19.05)

    make_title_slide(prs, data)
    print("  [1/8] 표지 슬라이드")
    make_overview_slide(prs, data)
    print("  [2/8] 프로젝트 개요")
    make_data_collection_slide(prs, data)
    print("  [3/8] 데이터 수집 현황")
    make_extraction_slide(prs, data)
    print("  [4/8] 추출 파이프라인")
    make_method_slide(prs, data)
    print("  [5/8] 합성 방법 분포")
    make_ml_slide(prs, data)
    print("  [6/8] ML 모델 결과")
    make_active_learning_slide(prs, data)
    print("  [7/8] 능동학습")
    make_inverse_design_slide(prs, data)
    print("  [8/8] 역설계")
    make_future_slide(prs, data)
    print("  [9/9] 향후 계획")
    make_summary_slide(prs, data)
    print("  [10/10] 요약")

    prs.save(PPT_PATH)
    print(f"\n✅  저장 완료: {PPT_PATH}")
    print(f"   슬라이드: 10장  |  16:9 와이드")


if __name__ == "__main__":
    main()
